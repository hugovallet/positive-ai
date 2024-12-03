from pathlib import Path
import logging
from datetime import datetime
from pathlib import Path
from typing import List, Optional

import click
from pptx import Presentation
from pydantic import BaseModel, EmailStr

from src import ROOT_DIR
from utils.click import SpecialHelpOrder
from utils.ppt import ExtendedSlide


class MemberInfo(BaseModel):
    member_name: str
    member_join_month: str
    member_logo_path: Optional[str]
    member_gatherer_firstame: str
    member_gatherer_lastname: str
    member_gatherer_email: EmailStr
    member_gatherer_picture: Optional[str]


class FirstPage(ExtendedSlide):
    """
    A class defining the main page.
    """

    def __init__(self, master_slide, member_info: MemberInfo):
        super().__init__(master_slide)
        self._member_info = member_info

    def fill(self):
        replace_text_in_shape(self.get_shape("Date Placeholder"), self._member_info.member_join_month)
        replace_text_in_shape(self.get_shape("Member Placeholder"), self._member_info.member_name)


class SecondPage(ExtendedSlide):

    def fill(self):
        """nothing to do for this one"""
        pass


class ThirdPage(ExtendedSlide):
    """
    A class defining the last page.
    """

    def __init__(self, master_slide, member_info: MemberInfo):
        super().__init__(master_slide)
        self._member_info = member_info

    def fill(self):
        combined_text = f""""
        {self._member_info.member_gatherer_firstame} {self._member_info.member_gatherer_lastname}
        référent Positive AI
        pour {self._member_info.member_name}
        {self._member_info.member_gatherer_email}
        """
        replace_text_in_shape(self.get_shape("Gatherer Info"), combined_text)


class MemberOnboardingDeck:
    """
    A class to fill all the slides of a template PowerPoint presentation.
    """

    def __init__(self,
                 infos: MemberInfo,
                 name: str,
                 template_path: Path = ROOT_DIR / 'templates' / '2024_09 Positive_AI_Flyer membres-template-fr.pptx'):
        self._log = logging.getLogger(__name__)
        self._template_path = Presentation(str(template_path))
        self._infos = infos

        # cached properties
        self._slides = None

    def get_layout(self, name: str):
        layouts = {layout.name: layout for layout in self._template_path.slide_layouts}
        return layouts[name]

    @property
    def slides(self) -> List[ExtendedSlide]:
        """
        Obtain a list of all ppt slides in the presentation (this property is cached).

        Notes:
            If you want to add slides, this is the method to modify!
        """
        if self._slides is None:
            slide_list = []

            # First slide
            front_layout = self.get_layout("first-page")
            master = self._template_path.slides.add_slide(front_layout)
            slide_list.append(FirstPage(master, member_info=self._infos))

            # Second slide
            disclaimer_layout = self.get_layout("second-page")
            master = self._template_path.slides.add_slide(disclaimer_layout)
            slide_list.append(SecondPage(master_slide=master))

            # Third slides
            end_layout = self.get_layout("third-page")
            master = self._template_path.slides.add_slide(end_layout)
            slide_list.append(ThirdPage(master_slide=master, member_info=self._infos))

            self._slides = slide_list

        return self._slides

    def save(self, file_path: Path = None):
        """
        Save in the current directory.
        """

        # if there is already a file in there, remove it
        if file_path.exists():
            file_path.unlink()

        # if the folder doesn't exist, create it
        if not file_path.parent.exists():
            file_path.parent.mkdir(True, True)

        # 2 fill all the slides with numbers and images
        for s in self.slides:
            s.fill()

        # save the underlying presentation object
        self._template_path.save(str(file_path))


def replace_text_in_shape(shape, new_text: str):
    """
    Replace the text in the Python ppt shape maintaining all formatting.
    """
    if shape.has_text_frame:
        if shape.is_placeholder:
            shape.text = new_text
        else:
            text_frame = shape.text_frame
            text_frame.paragraphs[0].runs[0].text = new_text
    else:
        raise TypeError("shape as no text box")


@click.group()
def cli():
    """List of commands"""


@click.command(
    help="Generate the employee starter presentation in english and french.",
    #help_priority=1
)
@click.option('--member-name',
              help="the name of the company joining positive AI",
              type=str,
              prompt=True
              )
@click.option('--member-join-month',
              help="the month the company joined positive AI",
              type=str,
              prompt=True
              )
@click.option('--member-gatherer-firstname',
              help="the firstname of the company gatherer",
              type=str,
              prompt=True
              )
@click.option('--member-gatherer-lastname',
              help="the lastname of the company gatherer",
              type=str,
              prompt=True
              )
@click.option('--member-gatherer-email',
              help="the email address of the company gatherer",
              type=str,
              prompt=True
              )
def generate(member_name, member_join_month, member_gatherer_firstname, member_gatherer_lastname, member_gatherer_email):
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")

    # Summarise member info from prompt
    infos = MemberInfo(
        member_name=member_name,
        member_join_month=member_join_month,
        member_gatherer_firstname=member_gatherer_firstname,
        member_gatherer_lastname=member_gatherer_lastname,
        member_gatherer_email=member_gatherer_email
    )

    # Build french deck
    fr_template_path = ROOT_DIR / "templates" / "2024_09 Positive_AI_Flyer membres-template-fr.pptx"
    french_deck = MemberOnboardingDeck(template_path=fr_template_path)
    filename = f"2024_09 Positive_AI_Flyer-{infos.member_name.lower()}-fr.pptx"
    french_deck.save(file_path=Path.cwd() / "positive-ai-generated" / "employee-onboarding" / filename)
