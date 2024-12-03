from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pydantic import BaseModel, EmailStr

from src.constants import ROOT_DIR
from utils.ppt import ExtendedSlide, replace_text_in_shape


class MemberInfo(BaseModel):
    member_name: str
    member_join_month: str
    member_logo_path: Optional[str] = None
    member_gatherer_firstname: str
    member_gatherer_lastname: str
    member_gatherer_email: EmailStr
    member_gatherer_photo_path: Optional[str] = None


class FirstPage(ExtendedSlide):
    """
    A class defining the main page.
    """

    def __init__(self, master_slide, member_info: MemberInfo):
        super().__init__(master_slide)
        self._member_info = member_info

    def fill(self):
        replace_text_in_shape(self.get_shape("Text Placeholder 1"), self._member_info.member_join_month)
        replace_text_in_shape(self.get_shape("Text Placeholder 2"), self._member_info.member_name)


class SecondPage(ExtendedSlide):

    def fill(self):
        """nothing to do for this one"""
        pass


class ThirdPage(ExtendedSlide):
    """
    A class defining the last page.
    """

    def __init__(self, master_slide, member_info: MemberInfo):
        super().__init__(master_slide)
        self._member_info = member_info

    def fill(self):
        combined_text = f"{self._member_info.member_gatherer_firstname} {self._member_info.member_gatherer_lastname}\nréférent Positive AI\npour {self._member_info.member_name}\n{self._member_info.member_gatherer_email}"
        replace_text_in_shape(self.get_shape("Text Placeholder 2"), combined_text)


class MemberOnboardingDeck:
    """
    A class to fill all the slides of a template PowerPoint presentation.
    """

    def __init__(self,
                 infos: MemberInfo,
                 template_path: Path = ROOT_DIR / 'templates' / '2024_09 Positive_AI_Flyer membres-template-fr.pptx'):
        self._template_path = Presentation(str(template_path))
        self._infos = infos

        # cached properties
        self._slides = None

    def get_layout(self, name: str):
        layouts = {layout.name: layout for layout in self._template_path.slide_layouts}
        return layouts[name]

    @property
    def slides(self) -> List[ExtendedSlide]:
        """
        Obtain a list of all ppt slides in the presentation (this property is cached).

        Notes:
            If you want to add slides, this is the method to modify!
        """
        if self._slides is None:
            slide_list = []

            # First slide
            front_layout = self.get_layout("first-page")
            master = self._template_path.slides.add_slide(front_layout)
            slide_list.append(FirstPage(master, member_info=self._infos))

            # Second slide
            disclaimer_layout = self.get_layout("second-page")
            master = self._template_path.slides.add_slide(disclaimer_layout)
            slide_list.append(SecondPage(master_slide=master))

            # Third slides
            end_layout = self.get_layout("third-page")
            master = self._template_path.slides.add_slide(end_layout)
            slide_list.append(ThirdPage(master_slide=master, member_info=self._infos))

            self._slides = slide_list

        return self._slides

    def save(self, file_path: Path = None):
        """
        Save in the current directory.
        """

        # if there is already a file in there, remove it
        if file_path.exists():
            file_path.unlink()

        # if the folder doesn't exist, create it
        if not file_path.parent.exists():
            file_path.parent.mkdir(exist_ok=True, parents=True, mode=0o770)

        # 2 fill all the slides with numbers and images
        for s in self.slides:
            s.fill()

        # save the underlying presentation object
        self._template_path.save(str(file_path))
