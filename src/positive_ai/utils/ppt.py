import abc
import logging
from pathlib import Path
from typing import Union
from pptx import Presentation
from pptx.shapes.placeholder import *
from pptx.slide import Slide

AnyPlaceholder = Union[
    LayoutPlaceholder,
    MasterPlaceholder,
    NotesSlidePlaceholder,
    SlidePlaceholder,
    ChartPlaceholder,
    PicturePlaceholder,
    PlaceholderGraphicFrame,
    PlaceholderPicture,
    TablePlaceholder,
]


class ExtendedSlide(Slide):
    """
    A class defining the slide holding all KPI data. It extends to base Slide class of pptx.
    """

    def __init__(self, master_slide: Slide, language: str):
        self._log = logging.getLogger(__name__)
        super().__init__(element=master_slide.element, part=master_slide.part)
        # cached prop
        self.__shape_name_to_index = None
        self._language = language

    @abc.abstractmethod
    def fill(self):
        """
        A method to define the set of action required to fill the slide with all contents (text, numbers, images, etc.)
        """
        pass

    @property
    def _shape_name_to_index(self):
        """
        A cached mapping between Shape names and shape indices.
        """
        if self.__shape_name_to_index is None:
            self.__shape_name_to_index = {
                s.name: index for index, s in enumerate(self.shapes)
            }
        return self.__shape_name_to_index

    def get_shape(self, shape_name) -> Union[Shape, AnyPlaceholder]:
        """
        A method to conveniently access Shape objects on the slide by their names.

        Args:
            shape_name: name of the shape

        Returns:
            shape: the Shape object
        """
        try:
            return self.shapes[self._shape_name_to_index[shape_name]]
        except KeyError:
            raise KeyError(
                f"Cannot find shape named {shape_name}. Available shapes: {list(self._shape_name_to_index.keys())}"
            )

    def __repr__(self) -> str:
        """
        This is the representation for our custom slide object. It displays all the contents of the slide to make it
        easier to understand / debug.
        """
        content = "\n".join(
            [
                f"shape={s.name}, type={type(s)}, index={index}, text='{s.text if hasattr(s, 'text') else ''}'"
                for index, s in enumerate(self.shapes)
            ]
        )
        return f"<{self.__class__.__name__}> at {hex(id(self))}\n{content}"


class Deck(object):
    __metaclass__ = abc.ABCMeta

    def __init__(self, infos, language: str, template_path: Path):
        self._template_path = Presentation(str(template_path))
        self._infos = infos
        self._language = language

        assert language in ("fr", "en"), f"Unsupported language '{language}'"

        # cached properties
        self._slides = None

    @property
    @abc.abstractmethod
    def slides(self):
        raise NotImplementedError(
            "You must override this attribute to define the slides of your presentation"
        )

    def get_layout(self, name: str):
        layouts = {layout.name: layout for layout in self._template_path.slide_layouts}
        return layouts[name]

    def save(self, file_path: Path = None):
        """Save in the provided directory."""

        # if there is already a file in there, remove it
        if file_path.exists():
            file_path.unlink()

        # if the folder doesn't exist, create it
        if not file_path.parent.exists():
            file_path.parent.mkdir(exist_ok=True, parents=True, mode=0o770)

        # 2 fill all the slides with numbers and images
        for s in self.slides:
            s.fill()

        # save the underlying presentation object
        self._template_path.save(str(file_path))


def replace_text_in_shape(shape: Shape, new_text: str):
    """
    Replace the text in the Python ppt shape maintaining all formatting.
    """
    if shape.has_text_frame:
        if shape.is_placeholder:
            shape.text = new_text
        else:
            text_frame = shape.text_frame
            text_frame.paragraphs[0].runs[0].text = new_text
    else:
        raise TypeError("shape as no text box")


def insert_image_in_shape(
    placeholder: Union[Shape, AnyPlaceholder],
    image_path: str,
    refit: bool = True,
    center: bool = False,
):
    # Get initial image placeholder left and top positions
    picture = placeholder.insert_picture(image_path)

    if refit:
        pos_left, pos_top = picture.left, picture.top
        available_width = picture.width
        available_height = picture.height

        image_width, image_height = picture.image.size
        placeholder_aspect_ratio = float(available_width) / float(available_height)
        image_aspect_ratio = float(image_width) / float(image_height)

        picture.crop_top = 0
        picture.crop_left = 0
        picture.crop_bottom = 0
        picture.crop_right = 0

        # ---if the placeholder is "wider" in aspect, shrink the picture width while
        # ---maintaining the image aspect ratio
        if placeholder_aspect_ratio > image_aspect_ratio:
            picture.width = int(image_aspect_ratio * available_height)
            picture.height = available_height

        # ---otherwise shrink the height
        else:
            picture.height = int(available_width / image_aspect_ratio)
            picture.width = available_width

        # Set the picture left and top position to the initial placeholder one
        picture.left, picture.top = pos_left, pos_top

        if center:
            picture.top = picture.top + int((available_height - picture.height) / 2)
            picture.left = picture.left + int((available_width - picture.width) / 2)
