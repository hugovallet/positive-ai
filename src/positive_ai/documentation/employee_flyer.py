from pathlib import Path
from typing import List

from pptx import Presentation

from positive_ai.documentation.data_model import MemberInfo
from positive_ai.utils.ppt import (
    ExtendedSlide,
    replace_text_in_shape,
    insert_image_in_shape,
    Deck,
)


class FirstPage(ExtendedSlide):
    """
    A class defining the main page.
    """

    def __init__(self, master_slide, member_info: MemberInfo, language: str):
        super().__init__(master_slide, language=language)
        self._member_info = member_info

    def fill(self):
        replace_text_in_shape(
            self.get_shape("Text Placeholder 1"), self._member_info.member_join_month
        )
        replace_text_in_shape(
            self.get_shape("Text Placeholder 2"), self._member_info.member_name
        )
        if self._member_info.member_logo_path:
            insert_image_in_shape(
                self.get_shape("Picture Placeholder 3"),
                self._member_info.member_logo_path,
                center=True,
            )


class SecondPage(ExtendedSlide):
    def fill(self):
        """nothing to do for this one"""
        pass


class ThirdPage(ExtendedSlide):
    """
    A class defining the last page.
    """

    def __init__(self, master_slide, member_info: MemberInfo, language: str):
        super().__init__(master_slide, language=language)
        self._member_info = member_info

    def fill(self):
        if self._language == "fr":
            combined_text = f"{self._member_info.member_gatherer_firstname} {self._member_info.member_gatherer_lastname}\nréférent Positive AI\npour {self._member_info.member_name}\n{self._member_info.member_gatherer_email}"
        elif self._language == "en":
            combined_text = f"{self._member_info.member_gatherer_firstname} {self._member_info.member_gatherer_lastname}\nPositive AI repr\nfor {self._member_info.member_name}\n{self._member_info.member_gatherer_email}"
        else:
            raise Exception(f"Unsupported language '{self._language}'")
        replace_text_in_shape(self.get_shape("Text Placeholder 2"), combined_text)

        if self._member_info.member_gatherer_photo_path:
            self.get_shape("Picture Placeholder 1").insert_picture(
                self._member_info.member_gatherer_photo_path
            )


class MemberOnboardingDeck(Deck):
    """
    A class to fill all the slides of a template PowerPoint presentation.
    """

    @property
    def slides(self) -> List[ExtendedSlide]:
        """
        Obtain a list of all ppt slides in the presentation (this property is cached).

        Notes:
            If you want to add slides, this is the method to modify!
        """
        if self._slides is None:
            slide_list = []

            # First slide
            front_layout = self.get_layout("first-page")
            master = self._template_path.slides.add_slide(front_layout)
            slide_list.append(
                FirstPage(master, member_info=self._infos, language=self._language)
            )

            # Second slide
            disclaimer_layout = self.get_layout("second-page")
            master = self._template_path.slides.add_slide(disclaimer_layout)
            slide_list.append(SecondPage(master_slide=master, language=self._language))

            # Third slides
            end_layout = self.get_layout("third-page")
            master = self._template_path.slides.add_slide(end_layout)
            slide_list.append(
                ThirdPage(
                    master_slide=master,
                    member_info=self._infos,
                    language=self._language,
                )
            )

            self._slides = slide_list

        return self._slides
