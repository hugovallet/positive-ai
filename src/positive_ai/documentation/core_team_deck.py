from datetime import datetime
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pydantic import BaseModel, EmailStr

from positive_ai.documentation.data_model import MemberInfo, CoreTeamMemberInfo
from positive_ai.utils.ppt import (
    ExtendedSlide,
    replace_text_in_shape,
    insert_image_in_shape,
    Deck,
)


def chunk_list(lst, size):
    return [lst[i : i + size] for i in range(0, len(lst), size)]


class FirstPage(ExtendedSlide):
    """
    A class defining the main page.
    """

    def __init__(self, master_slide, infos: MemberInfo, language: str):
        super().__init__(master_slide, language=language)
        self._infos = infos

    def fill(self):
        today = datetime.today().strftime("%b %d, %Y")
        replace_text_in_shape(self.get_shape("Subtitle 2"), today)
        if self._language == "fr":
            replace_text_in_shape(
                self.get_shape("Title 1"),
                "Conseil d'administration et Core Team Positive AI",
            )
        else:
            replace_text_in_shape(
                self.get_shape("Title 1"), "Positive AI Board and Core Team"
            )


class TrombiPage(ExtendedSlide):
    def __init__(self, master_slide, infos: List[CoreTeamMemberInfo], language: str):
        super().__init__(master_slide, language=language)
        self._infos = infos

    def set_title(self, title: str):
        replace_text_in_shape(self.get_shape("Title 1"), title)

    def fill(self):
        start_num = 1
        for i, member_info in enumerate(self._infos):
            shape_num = start_num + 4 * i
            insert_image_in_shape(
                self.get_shape(f"Picture Placeholder {shape_num + 1}"),
                member_info.ct_member_photo_path,
                refit=False,
            )
            replace_text_in_shape(
                self.get_shape(f"Text Placeholder {shape_num + 2}"),
                member_info.ct_member_firstname + " " + member_info.ct_member_lastname,
            )
            if self._language == "fr":
                replace_text_in_shape(
                    self.get_shape(f"Text Placeholder {shape_num + 3}"),
                    member_info.ct_member_title_fr,
                )
            else:
                replace_text_in_shape(
                    self.get_shape(f"Text Placeholder {shape_num + 3}"),
                    member_info.ct_member_title_en,
                )
            replace_text_in_shape(
                self.get_shape(f"Text Placeholder {shape_num + 4}"),
                member_info.ct_member_email,
            )

        # remove remaining placeholders
        for placeholder in self.shapes.placeholders:
            if placeholder.has_text_frame and placeholder.text_frame.text == "":
                sp = placeholder._sp
                sp.getparent().remove(sp)


class CoreTeamDeck(Deck):
    """
    A class to fill all the slides of a template PowerPoint presentation.
    """

    @property
    def slides(self) -> List[ExtendedSlide]:
        """
        Obtain a list of all ppt slides in the presentation (this property is cached).

        Notes:
            If you want to add slides, this is the method to modify!
        """
        if self._slides is None:
            slide_list = []
            # First slide
            layout = self.get_layout("Diapositive titre (lapis)")
            master = self._template_path.slides.add_slide(layout)
            slide_list.append(
                FirstPage(master, infos=self._infos, language=self._language)
            )
            board = [el for el in self._infos.all_members_info if el.ct_member_is_board]
            # create chunks of 8 members because trombi slide can handle only 8 members
            for chunk in chunk_list(board, 8):
                layout = self.get_layout("facebook-slide-dense")
                page = self._template_path.slides.add_slide(layout)
                page = TrombiPage(page, infos=chunk, language=self._language)
                if self._language == "fr":
                    page.set_title("Conseil d'administration Positive AI")
                else:
                    page.set_title("Positive AI board")
                slide_list.append(page)

            other = [
                el for el in self._infos.all_members_info if not el.ct_member_is_board
            ]
            # create chunks of 8 members because trombi slide can handle only 8 members
            for chunk in chunk_list(other, 8):
                layout = self.get_layout("facebook-slide-dense")
                page = self._template_path.slides.add_slide(layout)
                page = TrombiPage(page, infos=chunk, language=self._language)
                if self._language == "fr":
                    page.set_title("Core Team Positive AI")
                else:
                    page.set_title("Positive AI Core Team")
                slide_list.append(page)

            self._slides = slide_list

        return self._slides
