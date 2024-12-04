from datetime import datetime
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pydantic import BaseModel, EmailStr

from positive_ai.documentation.data_model import MemberInfo
from positive_ai.utils.ppt import (
    ExtendedSlide,
    replace_text_in_shape,
    insert_image_in_shape,
    Deck,
)


def chunk_list(lst, size):
    return [lst[i : i + size] for i in range(0, len(lst), size)]


class FirstPage(ExtendedSlide):
    """
    A class defining the main page.
    """

    def __init__(self, master_slide, infos: MemberInfo, language: str):
        super().__init__(master_slide, language=language)
        self._infos = infos

    def fill(self):
        today = datetime.today().strftime("%b %d, %Y")
        replace_text_in_shape(self.get_shape("Subtitle 2"), today)
        if self._language == "fr":
            replace_text_in_shape(self.get_shape("Title 1"), "Communauté Positive AI")
        else:
            replace_text_in_shape(self.get_shape("Title 1"), "Positive AI Community")


class TrombiPage(ExtendedSlide):
    def __init__(self, master_slide, infos: List[MemberInfo], language: str):
        super().__init__(master_slide, language=language)
        self._infos = infos

    def fill(self):
        if self._language == "fr":
            replace_text_in_shape(
                self.get_shape("Title 1"), "Communauté PAI - Référents Entreprise"
            )
        else:
            replace_text_in_shape(
                self.get_shape("Title 1"), "PAI Community - Gatherers"
            )

        start_num = 2
        for i, member_info in enumerate(self._infos):
            shape_num = start_num + 6 * i
            insert_image_in_shape(
                self.get_shape(f"Picture Placeholder {shape_num}"),
                member_info.member_logo_path,
                refit=True,
                center=True,
            )
            insert_image_in_shape(
                self.get_shape(f"Picture Placeholder {shape_num + 1}"),
                member_info.member_gatherer_photo_path,
                refit=False,
            )
            replace_text_in_shape(
                self.get_shape(f"Text Placeholder {shape_num + 2}"),
                member_info.member_gatherer_firstname
                + " "
                + member_info.member_gatherer_lastname,
            )
            if self._language == "fr":
                replace_text_in_shape(
                    self.get_shape(f"Text Placeholder {shape_num + 3}"),
                    member_info.member_gatherer_title_fr,
                )
            else:
                replace_text_in_shape(
                    self.get_shape(f"Text Placeholder {shape_num + 3}"),
                    member_info.member_gatherer_title_en,
                )
            replace_text_in_shape(
                self.get_shape(f"Text Placeholder {shape_num + 4}"),
                member_info.member_gatherer_email,
            )
            if self._language == "fr":
                replace_text_in_shape(
                    self.get_shape(f"Text Placeholder {shape_num + 5}"),
                    member_info.member_gatherer_desc_fr,
                )
            else:
                replace_text_in_shape(
                    self.get_shape(f"Text Placeholder {shape_num + 5}"),
                    member_info.member_gatherer_desc_en,
                )

        # remove remaining placeholders
        for placeholder in self.shapes.placeholders:
            if placeholder.has_text_frame and placeholder.text_frame.text == "":
                sp = placeholder._sp
                sp.getparent().remove(sp)


class CommunityDeck(Deck):
    """
    A class to fill all the slides of a template PowerPoint presentation.
    """

    @property
    def slides(self) -> List[ExtendedSlide]:
        """
        Obtain a list of all ppt slides in the presentation (this property is cached).

        Notes:
            If you want to add slides, this is the method to modify!
        """
        if self._slides is None:
            slide_list = []
            # First slide
            layout = self.get_layout("Diapositive titre (lapis)")
            master = self._template_path.slides.add_slide(layout)
            slide_list.append(
                FirstPage(master, infos=self._infos, language=self._language)
            )
            # create chunks of 4 members because trombi slide can handle only 4 members
            for chunk in chunk_list(self._infos.all_members_info, 4):
                layout = self.get_layout("trombi-slide")
                page = self._template_path.slides.add_slide(layout)
                slide_list.append(
                    TrombiPage(page, infos=chunk, language=self._language)
                )

            self._slides = slide_list

        return self._slides
