import abc
import logging
from typing import Union

from pptx.shapes.placeholder import *
from pptx.slide import Slide

AnyPlaceholder = Union[LayoutPlaceholder, MasterPlaceholder, NotesSlidePlaceholder, SlidePlaceholder, ChartPlaceholder,
                       PicturePlaceholder, PlaceholderGraphicFrame, PlaceholderPicture, TablePlaceholder]


class ExtendedSlide(Slide):
    """
    A class defining the slide holding all KPI data. It extends to base Slide class of pptx.
    """

    def __init__(self, master_slide: Slide, language: str):
        self._log = logging.getLogger(__name__)
        super().__init__(element=master_slide.element, part=master_slide.part)
        # cached prop
        self.__shape_name_to_index = None
        self._language = language

    @abc.abstractmethod
    def fill(self):
        """
        A method to define the set of action required to fill the slide with all contents (text, numbers, images, etc.)
        """
        pass

    @property
    def _shape_name_to_index(self):
        """
        A cached mapping between Shape names and shape indices.
        """
        if self.__shape_name_to_index is None:
            self.__shape_name_to_index = {s.name: index for index, s in enumerate(self.shapes)}
        return self.__shape_name_to_index

    def get_shape(self, shape_name) -> Union[Shape, AnyPlaceholder]:
        """
        A method to conveniently access Shape objects on the slide by their names.

        Args:
            shape_name: name of the shape

        Returns:
            shape: the Shape object
        """
        try:
            return self.shapes[self._shape_name_to_index[shape_name]]
        except KeyError:
            raise KeyError(f"Cannot find shape named {shape_name}. Available shapes: {list(self._shape_name_to_index.keys())}")

    def __repr__(self) -> str:
        """
        This is the representation for our custom slide object. It displays all the contents of the slide to make it
        easier to understand / debug.
        """
        content = "\n".join(
            [f"shape={s.name}, type={type(s)}, index={index}, text='{s.text if hasattr(s, 'text') else ''}'" for index, s
             in enumerate(self.shapes)])
        return f"<{self.__class__.__name__}> at {hex(id(self))}\n{content}"


def replace_text_in_shape(shape, new_text: str):
    """
    Replace the text in the Python ppt shape maintaining all formatting.
    """
    if shape.has_text_frame:
        if shape.is_placeholder:
            shape.text = new_text
        else:
            text_frame = shape.text_frame
            text_frame.paragraphs[0].runs[0].text = new_text
    else:
        raise TypeError("shape as no text box")
